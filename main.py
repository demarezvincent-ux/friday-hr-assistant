import streamlit as st
from supabase import create_client, Client
import requests
import pdfplumber
import docx
import re
import time
import os
import datetime
import uuid
import logging
import hashlib
import json
import pandas as pd
from pptx import Presentation
from huggingface_hub import InferenceClient
from services.rag_controller import get_context_with_strategy

# Setup logging for debugging
logger = logging.getLogger(__name__)
logging.basicConfig(level=logging.INFO)

# --- 1. CONFIGURATION & PREMIUM STYLING ---
st.set_page_config(page_title="FRIDAY", page_icon="⚡", layout="wide")

st.markdown("""
<style>
    /* ============================================
       FRIDAY - Premium Design System v3.0
       ============================================ */
    @import url('https://fonts.googleapis.com/css2?family=Playfair+Display:ital,wght@0,400;0,600;0,700;1,400&family=Inter:wght@300;400;500;600&display=swap');

    :root {
        --primary: #1A3C34;
        --primary-light: #2A5248;
        --background: #F9F9F7;
        --sidebar-bg: #F0F0EE;
        --text-primary: #1A3C34;
        --text-secondary: #5C6F68;
        --white: #FFFFFF;
        --border: #E6E6E3;
    }

    /* Global Reset & Typography */
    html, body, [class*="css"] {
        font-family: 'Inter', sans-serif;
        color: var(--text-primary);
        background-color: var(--background);
    }
    
    h1, h2, h3, h4, h5, h6, .playfair {
        font-family: 'Playfair Display', serif !important;
        color: var(--text-primary) !important;
    }

    /* Streamlit App Container */
    .stApp {
        background-color: var(--background);
    }

    /* ============================================
       SIDEBAR
       ============================================ */
    [data-testid="stSidebar"] {
        background-color: var(--sidebar-bg);
        border-right: 1px solid rgba(26, 60, 52, 0.06);
    }
    
    [data-testid="stSidebar"] hr {
        margin: 24px 0;
        border-color: rgba(26, 60, 52, 0.1) !important;
    }

    /* Sidebar Buttons */
    [data-testid="stSidebar"] .stButton > button {
        background-color: transparent;
        color: var(--text-secondary);
        border: none;
        text-align: left;
        padding-left: 12px;
        font-weight: 500;
        transition: all 0.2s ease;
        justify-content: flex-start;
    }
    
    [data-testid="stSidebar"] .stButton > button:hover {
        color: var(--primary);
        background-color: rgba(26, 60, 52, 0.04);
        transform: translateX(4px);
    }

    /* Selected state for Sidebar Buttons handled by 'type="primary"' in Python 
       but let's override the primary style to fit the theme */
    [data-testid="stSidebar"] .stButton > button[kind="primary"] {
        background-color: var(--primary) !important;
        color: var(--white) !important;
        box-shadow: 0 4px 12px rgba(26, 60, 52, 0.15);
        border-radius: 8px;
    }

    /* ============================================
       MAIN CONTENT & CHAT
       ============================================ */
    
    /* Input Fields */
    .stTextInput > div > div > input, 
    .stChatInput > div > div > textarea {
        background-color: var(--white);
        border: 1px solid var(--border);
        border-radius: 12px;
        color: var(--text-primary);
        box-shadow: 0 2px 8px rgba(0,0,0,0.02);
    }
    
    .stTextInput > div > div > input:focus,
    .stChatInput > div > div > textarea:focus {
        border-color: var(--primary);
        box-shadow: 0 0 0 2px rgba(26, 60, 52, 0.1);
    }

    /* Chat Messages */
    [data-testid="stChatMessage"] {
        background-color: transparent;
        gap: 1.5rem;
    }

    /* User Message */
    [data-testid="stChatMessage"][data-testid="user"] {
        background-color: transparent; 
    }
    
    /* Assistant Message */
    [data-testid="stChatMessage"][data-testid="assistant"] {
        background-color: transparent;
    }

    /* Avatars */
    div[data-testid="stChatMessage"] div[data-testid="stColorBlock"],
    [data-testid="chatAvatarIcon-user"], 
    [data-testid="chatAvatarIcon-assistant"] {
        background-color: var(--primary) !important;
        border-radius: 50%;
        color: white !important;
    }
    
    div[data-testid="stChatMessage"] img[data-testid="stChatAvatar"] {
        border-radius: 50%;
        border: 2px solid var(--primary);
    }

    /* ============================================
       FILE LIST & ICONS
       ============================================ */
    .file-list-card {
        background: var(--white);
        border-radius: 12px;
        padding: 16px;
        margin-bottom: 12px;
        border: 1px solid var(--border);
        box-shadow: 0 2px 6px rgba(0,0,0,0.02);
    }

    /* Center align control buttons - targeting the Documents page columns */
    /* This targets buttons inside the columns used for the file list */
    [data-testid="stVerticalBlock"] [data-testid="column"] button {
        margin: 0 auto;
        display: block;
    }

    /* Status Dot */
    .status-dot {
        width: 8px;
        height: 8px;
        border-radius: 50%;
        background-color: var(--primary);
    }

    /* ============================================
       ANIMATIONS
       ============================================ */
    @keyframes fadeIn { from { opacity: 0; transform: translateY(10px); } to { opacity: 1; transform: translateY(0); } }
    
    .element-container, .stChatMessage {
        animation: fadeIn 0.4s ease-out;
    }

    .logo-animated {
        animation: fadeIn 1s ease-out;
    }

    /* Thinking Spinner */
    .thinking-spinner { display: flex; gap: 6px; padding: 12px 0; align-items: center; }
    .thinking-spinner span:not(.thinking-text) {
        width: 6px; height: 6px; background: var(--primary);
        border-radius: 50%; opacity: 0.6;
        animation: pulse 1.4s infinite ease-in-out;
    }
    .thinking-spinner span:nth-child(1) { animation-delay: -0.32s; }
    .thinking-spinner span:nth-child(2) { animation-delay: -0.16s; }
    @keyframes pulse { 0%, 80%, 100% { transform: scale(0); } 40% { transform: scale(1); } }
    
    /* Scrollbar */
    ::-webkit-scrollbar { width: 6px; }
    ::-webkit-scrollbar-track { background: transparent; }
    ::-webkit-scrollbar-thumb { background: #D2D2D7; border-radius: 10px; }
    ::-webkit-scrollbar-thumb:hover { background: #86868B; }

    /* Source Citation Styles */
    .sources-section {
        margin-top: 16px;
        padding: 12px 16px;
        background: rgba(26, 60, 52, 0.03);
        border-radius: 12px;
        border: 1px solid rgba(26, 60, 52, 0.08);
    }
    
    .sources-section details {
        cursor: pointer;
    }
    
    .sources-section summary {
        font-size: 13px;
        font-weight: 500;
        color: #5C6F68;
        list-style: none;
        display: flex;
        align-items: center;
        gap: 8px;
    }
    
    .sources-section summary::-webkit-details-marker { display: none; }
    
    .sources-section summary::before {
        content: '▶';
        font-size: 10px;
        transition: transform 0.2s ease;
    }
    
    .sources-section details[open] summary::before {
        transform: rotate(90deg);
    }
    
    .sources-list {
        display: flex;
        flex-wrap: wrap;
        gap: 8px;
        margin-top: 12px;
        padding-left: 18px;
    }
    
    .source-pill {
        padding: 6px 12px;
        border-radius: 16px;
        font-size: 12px;
        font-weight: 500;
        display: inline-flex;
        align-items: center;
        gap: 6px;
    }
    
    .legal-source-pill {
        background-color: rgba(16, 124, 65, 0.10);
        color: #0d6b3a;
        border: 1px solid rgba(16, 124, 65, 0.18);
    }
    
    .sector-source-pill {
        background-color: rgba(217, 119, 6, 0.10);
        color: #92400e;
        border: 1px solid rgba(217, 119, 6, 0.18);
    }
    
    .company-source-pill {
        background-color: rgba(37, 99, 235, 0.08);
        color: #1e40af;
        border: 1px solid rgba(37, 99, 235, 0.15);
    }
    
    .source-group-header {
        font-size: 11px;
        font-weight: 600;
        text-transform: uppercase;
        letter-spacing: 0.5px;
        margin-bottom: 6px;
        margin-top: 8px;
    }
    
    .source-group-header.legal { color: #0d6b3a; }
    .source-group-header.sector { color: #92400e; }
    .source-group-header.company { color: #1e40af; }
    
    .no-sources {
        font-size: 13px;
        color: #86868B;
        font-style: italic;
        display: flex;
        align-items: center;
        gap: 8px;
    }
    
    .single-source {
        font-size: 13px;
        color: #5C6F68;
        display: flex;
        align-items: center;
        gap: 8px;
    }
    
    .single-source .source-name {
        font-weight: 500;
        color: #1A3C34;
    }

    .onboarding-shell {
        background: linear-gradient(135deg, #ffffff 0%, #f5f7f4 100%);
        border: 1px solid rgba(26, 60, 52, 0.1);
        border-radius: 18px;
        padding: 28px;
        box-shadow: 0 18px 50px rgba(26, 60, 52, 0.08);
        margin-bottom: 18px;
    }

    .onboarding-title {
        font-family: 'Playfair Display', serif;
        font-size: 42px;
        line-height: 1.1;
        margin: 0 0 10px 0;
    }

    .onboarding-subtitle {
        color: #5c6f68;
        margin-bottom: 20px;
    }

    .legal-upload-card {
        background: #ffffff;
        border: 1px solid rgba(26, 60, 52, 0.12);
        border-radius: 14px;
        padding: 18px;
        margin-bottom: 16px;
    }

    @keyframes softFloatIn {
        from { opacity: 0; transform: translateY(10px); }
        to { opacity: 1; transform: translateY(0); }
    }

    .animated-panel {
        animation: softFloatIn 0.35s ease-out;
    }

    .onboarding-step-card {
        background: #ffffff;
        border: 1px solid rgba(26, 60, 52, 0.10);
        border-radius: 14px;
        padding: 18px;
        box-shadow: 0 8px 24px rgba(26, 60, 52, 0.06);
        transition: transform 0.25s ease, box-shadow 0.25s ease;
    }

    .onboarding-step-card:hover {
        transform: translateY(-2px);
        box-shadow: 0 12px 28px rgba(26, 60, 52, 0.09);
    }

    .stButton > button {
        transition: transform 0.2s ease, box-shadow 0.2s ease;
    }

    .stButton > button:hover {
        transform: translateY(-1px);
    }

    [data-testid="stProgressBar"] > div > div > div > div {
        transition: width 0.35s ease;
    }

    @keyframes fadeWideIn {
        from { opacity: 0; transform: scale(0.985); }
        to { opacity: 1; transform: scale(1); }
    }

    @keyframes typing {
        from { width: 0; }
        to { width: 100%; }
    }

    @keyframes blinkCaret {
        50% { border-color: transparent; }
    }

    .welcome-shell {
        display: flex;
        justify-content: center;
        align-items: center;
        min-height: 48vh;
        animation: fadeWideIn 0.6s ease-out;
    }

    .welcome-typewriter {
        overflow: hidden;
        white-space: nowrap;
        border-right: 2px solid var(--primary);
        font-family: 'Playfair Display', serif;
        font-size: clamp(36px, 5vw, 68px);
        color: var(--primary);
        letter-spacing: 0.02em;
        width: 0;
        animation: typing 2.2s steps(18, end) forwards, blinkCaret 0.85s step-end infinite;
    }

</style>
""", unsafe_allow_html=True)


# --- SECRETS HANDLING ---
def get_secret(key_name):
    """Get secret from environment or Streamlit secrets with validation."""
    if not key_name or not isinstance(key_name, str):
        return None
    
    # Try environment variable first
    if key_name in os.environ:
        value = os.environ[key_name]
        if isinstance(value, str) and value.strip():
            return value.strip()
    
    # Try Streamlit secrets
    try:
        if key_name in st.secrets:
            value = st.secrets[key_name]
            if isinstance(value, str) and value.strip():
                return value.strip()
    except Exception as e:
        logger.warning(f"Error accessing Streamlit secrets for {key_name}: {e}")
    
    return None

SUPABASE_URL = get_secret("SUPABASE_URL")
SUPABASE_KEY = get_secret("SUPABASE_KEY")
FIXED_GROQ_KEY = get_secret("FIXED_GROQ_KEY")
HF_API_KEY = get_secret("HF_API_KEY")

# Validate API keys are strings and not empty
api_keys_valid = all([
    isinstance(SUPABASE_URL, str) and SUPABASE_URL.strip(),
    isinstance(SUPABASE_KEY, str) and SUPABASE_KEY.strip(),
    isinstance(FIXED_GROQ_KEY, str) and FIXED_GROQ_KEY.strip(),
    isinstance(HF_API_KEY, str) and HF_API_KEY.strip()
])

if not api_keys_valid:
    st.error("❌ Missing or invalid API Keys. Please check your Secrets or Environment Variables.")
    st.stop()

# --- 2. STATE ---
if "authenticated" not in st.session_state: st.session_state.authenticated = False
if "company_id" not in st.session_state: st.session_state.company_id = None
if "current_chat_id" not in st.session_state: st.session_state.current_chat_id = None
if "view" not in st.session_state: st.session_state.view = "chat"
if "onboarding_required" not in st.session_state: st.session_state.onboarding_required = False
if "onboarding_step" not in st.session_state: st.session_state.onboarding_step = 1
if "show_welcome_anim" not in st.session_state: st.session_state.show_welcome_anim = False

@st.cache_resource
def init_supabase():
    if not SUPABASE_URL or not SUPABASE_KEY:
        raise ValueError("Supabase URL and key are required")
    return create_client(SUPABASE_URL, SUPABASE_KEY)

supabase = init_supabase()

# --- 3. BACKEND LOGIC ---

from services.agentic.rate_limiter import get_huggingface_limiter

def get_embeddings_batch(texts):
    model_id = "sentence-transformers/paraphrase-multilingual-MiniLM-L12-v2"
    client = InferenceClient(token=HF_API_KEY)
    clean_texts = [t.replace("\n", " ").strip() for t in texts]
    backoff_times = [2, 4, 8, 16]
    
    # Rate limit to prevent 429 errors
    hf_limiter = get_huggingface_limiter()
    hf_limiter.wait_if_needed()

    for wait_time in backoff_times:
        try:
            embeddings = client.feature_extraction(clean_texts, model=model_id)
            if hasattr(embeddings, "tolist"): return embeddings.tolist()
            return embeddings
        except:
            time.sleep(wait_time)
    return None

def sanitize_filename(filename):
    if not filename or not isinstance(filename, str):
        return "unknown_file"
    # Limit filename length to prevent path traversal
    name = filename[:100].replace(" ", "_")
    # Remove path traversal attempts and dangerous characters
    name = re.sub(r'[/\\:*?"<>|]', '', name)
    name = re.sub(r'\.\.', '', name)  # Remove directory traversal
    name = re.sub(r'[^a-zA-Z0-9._-]', '', name)
    # Ensure filename doesn't start with a dot (hidden files)
    if name.startswith('.'):
        name = 'file_' + name[1:] if len(name) > 1 else 'file'
    return name or "unknown_file"

def normalize_query(query):
    """Normalize query but keep original intent."""
    return query.strip().lower()

def extract_text_from_pdf(file):
    text = ""
    try:
        with pdfplumber.open(file) as pdf:
            for page in pdf.pages:
                tables = page.extract_tables()
                if tables:
                    for table in tables:
                        clean_table = [[(str(cell) if cell else "").replace("\n", " ") for cell in row] for row in table]
                        if clean_table:
                            try:
                                header = "| " + " | ".join(clean_table[0]) + " |"
                                sep = "| " + " | ".join(["---"] * len(clean_table[0])) + " |"
                                body = "\n".join(["| " + " | ".join(row) + " |" for row in clean_table[1:]])
                                text += f"\n{header}\n{sep}\n{body}\n\n"
                            except: pass
                page_text = page.extract_text()
                if page_text: text += page_text + "\n"
    except: return ""
    return text

def extract_text_from_excel(file) -> str:
    """Extract text from .xlsx files using pandas."""
    try:
        file.seek(0)
        dfs = pd.read_excel(file, sheet_name=None, engine='openpyxl')
        text_parts = []
        for sheet_name, df in dfs.items():
            text_parts.append(f"## Sheet: {sheet_name}\n")
            if not df.empty:
                text_parts.append(df.to_markdown(index=False))
        return "\n\n".join(text_parts)
    except Exception as e:
        logger.warning(f"Excel extraction failed: {e}")
        return ""

def extract_text_from_pptx(file) -> str:
    """Extract text from .pptx files slides."""
    try:
        file.seek(0)
        prs = Presentation(file)
        text_parts = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = [f"## Slide {slide_num}"]
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            text_parts.append("\n".join(slide_text))
        return "\n\n".join(text_parts)
    except Exception as e:
        logger.warning(f"PPTX extraction failed: {e}")
        return ""

def extract_file_metadata(file) -> dict:
    """Extract metadata from file properties."""
    meta = {"title": file.name, "created_date": None, "author": None}
    try:
        file.seek(0)
        fname_lower = file.name.lower()
        if fname_lower.endswith(".docx"):
            doc = docx.Document(file)
            props = doc.core_properties
            meta["author"] = props.author if props.author else None
            meta["title"] = props.title if props.title else file.name
        elif fname_lower.endswith(".pptx"):
            prs = Presentation(file)
            props = prs.core_properties
            meta["author"] = props.author if props.author else None
            meta["title"] = props.title if props.title else file.name
        elif fname_lower.endswith(".pdf"):
            with pdfplumber.open(file) as pdf:
                md = pdf.metadata
                if md:
                    meta["author"] = md.get("Author")
                    meta["title"] = md.get("Title") if md.get("Title") else file.name
    except: pass
    file.seek(0)
    return meta

def extract_text_and_metadata(file):
    """Extract text/metadata from supported file types and return (text, metadata, ext)."""
    text = ""
    file_metadata = {"title": sanitize_filename(file.name)}
    ext = file.name.lower().split('.')[-1]
    if ext == "pdf":
        file_metadata = extract_file_metadata(file)
        text = extract_text_from_pdf(file)
    elif ext == "docx":
        file_metadata = extract_file_metadata(file)
        doc = docx.Document(file)
        text = "\n".join([p.text for p in doc.paragraphs])
    elif ext == "xlsx":
        text = extract_text_from_excel(file)
    elif ext == "pptx":
        file_metadata = extract_file_metadata(file)
        text = extract_text_from_pptx(file)
    else:
        return "", file_metadata, ext
    return text, file_metadata, ext

def get_mime_type_for_extension(ext: str) -> str:
    if ext == "pdf":
        return "application/pdf"
    if ext == "docx":
        return "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    if ext == "xlsx":
        return "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    if ext == "pptx":
        return "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    return "application/octet-stream"

def recursive_chunk_text(text: str, chunk_size: int = 800, overlap: int = 150) -> list:
    """Recursive Character Splitter for semantic preservation."""
    if not text: return []
    paragraphs = text.split("\n\n")
    chunks = []
    current_chunk = ""
    for para in paragraphs:
        para = para.strip()
        if not para: continue
        if len(current_chunk) + len(para) + 2 < chunk_size:
            current_chunk += para + "\n\n"
        else:
            if current_chunk: chunks.append(current_chunk.strip())
            if len(para) > chunk_size:
                parts = [para[i:i+chunk_size] for i in range(0, len(para), chunk_size-overlap)]
                chunks.extend(parts)
                current_chunk = ""
            else:
                current_chunk = para + "\n\n"
    if current_chunk.strip(): chunks.append(current_chunk.strip())
    return chunks

# Legacy alias
def smart_chunking(text, chunk_size=500, overlap=100):
    return recursive_chunk_text(text, chunk_size, overlap)

# --- DATABASE OPERATIONS ---

def register_document(filename, company_id, metadata=None):
    try:
        doc_data = {"company_id": company_id, "filename": filename, "is_active": True}
        if metadata:
            if metadata.get("title"):
                doc_data["title"] = metadata.get("title")
            if metadata.get("author"):
                doc_data["author"] = metadata.get("author")
        result = supabase.table("documents").insert(doc_data).execute()
        logger.info(f"Document registered: {filename}, result: {result.data}")
        return True
    except Exception as e:
        logger.error(f"Failed to register document {filename}: {e}")
        return False

def check_if_document_exists(filename, company_id):
    try:
        res = supabase.table("documents").select("id").eq("company_id", company_id).eq("filename", filename).execute()
        return len(res.data) > 0
    except: return False

def get_all_documents(company_id):
    try:
        result = supabase.table("documents").select("*").eq("company_id", company_id).order('is_active', desc=True).order('created_at', desc=True).execute()
        logger.info(f"Fetched {len(result.data)} documents for company {company_id}")
        return result.data
    except Exception as e:
        logger.error(f"Failed to fetch documents: {e}")
        return []

def get_chunk_counts(company_id):
    """Get the number of chunks stored for each document in a company.
    
    Strategy: Query chunks by filename (ASCII) to avoid emoji encoding issues.
    We get the list of document filenames from the documents table and check each one.
    """
    if not company_id or not isinstance(company_id, str):
        logger.error("Invalid company_id provided to get_chunk_counts")
        return {}
    
    try:
        # First get the list of filenames we need to check
        docs = supabase.table("documents").select("filename").eq("company_id", company_id).execute()
        if not docs or not docs.data:
            return {}
        
        filenames = [doc['filename'] for doc in docs.data if isinstance(doc, dict) and 'filename' in doc]
        
        chunk_counts = {}
        for filename in filenames:
            if not filename or not isinstance(filename, str):
                continue
            # Query chunks by filename (ASCII) - no emoji encoding issues
            result = supabase.table("document_chunks").select("id", count="exact").eq(
                "metadata->>filename", filename
            ).execute()
            chunk_counts[filename] = result.count if hasattr(result, 'count') and result.count else 0
        
        logger.info(f"Chunk counts for {company_id}: {chunk_counts}")
        return chunk_counts
    except Exception as e:
        logger.error(f"Failed to get chunk counts: {e}")
        return {}





def toggle_document_status(filename, company_id, current_status):
    try:
        supabase.table("documents").update({"is_active": not current_status}).eq("company_id", company_id).eq("filename", filename).execute()
        return True
    except: return False

def delete_document(filename, company_id):
    try:
        supabase.table("documents").delete().eq("company_id", company_id).eq("filename", filename).execute()
        supabase.table("document_chunks").delete().eq("metadata->>company_id", company_id).eq("metadata->>filename", filename).execute()
        try: supabase.storage.from_("documents").remove([f"{company_id}/{filename}"])
        except: pass
        return True
    except: return False

COMPANY_PROFILE_FILENAME = "__company_profile__.md"
COMPANY_PROFILE_SNAPSHOT_TYPE = "company_profile_snapshot"
REQUIRED_PROFILE_FIELDS = [
    "company_name",
    "sector",
    "joint_committees",
    "operations",
    "headquarters",
    "countries",
    "employees_total",
    "employees_belgium",
    "contract_types",
    "weekly_hours",
    "payroll_frequency_blue",
    "payroll_frequency_white",
    "existing_policies",
    "priorities",
]

SECTOR_PC_SUGGESTIONS = {
    "food": ["PC 118", "PC 119", "PC 220"],
    "hospitality": ["PC 302"],
    "horeca": ["PC 302"],
    "retail": ["PC 201", "PC 311", "PC 312"],
    "manufacturing": ["PC 111", "PC 124", "PC 200"],
    "construction": ["PC 124"],
    "logistics": ["PC 140.03", "PC 226"],
    "transport": ["PC 140.03"],
    "healthcare": ["PC 330", "PC 331"],
    "services": ["PC 200"],
}

INDUSTRY_PC_SUGGESTIONS = {
    "Hospitality": ["PC 302"],
    "Manufacturing": ["PC 111", "PC 124", "PC 200"],
    "Retail": ["PC 201", "PC 311", "PC 312"],
    "Logistics/Transport": ["PC 140.03", "PC 226"],
    "Healthcare": ["PC 330", "PC 331"],
    "Professional Services": ["PC 200", "PC 218"],
    "Other": ["PC 200"],
}

COMMON_PC_OPTIONS = [
    "PC 100", "PC 109", "PC 111", "PC 116", "PC 118", "PC 119", "PC 124",
    "PC 140.03", "PC 145", "PC 149.01", "PC 200", "PC 201", "PC 202", "PC 207",
    "PC 209", "PC 218", "PC 220", "PC 226", "PC 302", "PC 311", "PC 312",
    "PC 322", "PC 327", "PC 330", "PC 331", "PC 337",
]

def suggest_joint_committees(sector: str, industry_cluster: str) -> list:
    sector_text = (sector or "").lower()
    suggestions = set(INDUSTRY_PC_SUGGESTIONS.get(industry_cluster, []))
    for keyword, pcs in SECTOR_PC_SUGGESTIONS.items():
        if keyword in sector_text:
            suggestions.update(pcs)
    if not suggestions:
        suggestions.add("PC 200")
    return sorted(suggestions)

def normalize_profile_snapshot_keys(snapshot: dict) -> dict:
    if not isinstance(snapshot, dict):
        return {}
    normalized = dict(snapshot)
    # Legacy key migrations
    if normalized.get("joint_committee") and not normalized.get("joint_committees"):
        normalized["joint_committees"] = normalized.get("joint_committee")
    if normalized.get("payroll_frequency"):
        if not normalized.get("payroll_frequency_blue"):
            normalized["payroll_frequency_blue"] = normalized.get("payroll_frequency")
        if not normalized.get("payroll_frequency_white"):
            normalized["payroll_frequency_white"] = normalized.get("payroll_frequency")
    # If legacy snapshot has all core fields, mark as complete.
    legacy_core = [
        "company_name", "sector", "operations", "headquarters", "countries",
        "employees_total", "employees_belgium", "contract_types", "weekly_hours",
        "existing_policies", "priorities"
    ]
    if not normalized.get("_profile_completed"):
        if all(str(normalized.get(k, "")).strip() for k in legacy_core):
            normalized["_profile_completed"] = True
    return normalized

def company_profile_exists(company_id):
    return check_if_document_exists(COMPANY_PROFILE_FILENAME, company_id)

def get_company_profile_snapshot(company_id):
    """Read the indexed profile back from document chunks and parse key/value pairs."""
    try:
        # Read the latest snapshot for THIS company only.
        snapshot_res = supabase.table("document_chunks").select("id, content").eq(
            "metadata->>company_id", company_id
        ).eq(
            "metadata->>document_type", COMPANY_PROFILE_SNAPSHOT_TYPE
        ).order("id", desc=True).limit(1).execute()
        if snapshot_res.data:
            raw = (snapshot_res.data[0] or {}).get("content", "")
            if raw.startswith("__PROFILE_JSON__:"):
                payload = raw.replace("__PROFILE_JSON__:", "", 1).strip()
                parsed = {k: str(v) for k, v in json.loads(payload).items() if v is not None}
                return normalize_profile_snapshot_keys(parsed)

        # Fallback to parsing regular profile chunks
        result = supabase.table("document_chunks").select("content").eq(
            "metadata->>company_id", company_id
        ).eq(
            "metadata->>filename", COMPANY_PROFILE_FILENAME
        ).execute()
        if not result.data:
            return {}

        combined = "\n".join([(row.get("content") or "") for row in result.data if isinstance(row, dict)])
        snapshot = {}
        alias_map = {
            "joint_committee_pc": "joint_committees",
            "joint_committees_pc": "joint_committees",
            "countries_of_operation": "countries",
            "total_employees": "employees_total",
            "employees_in_belgium": "employees_belgium",
            "standard_weekly_hours": "weekly_hours",
            "payroll_frequency_bluecollar": "payroll_frequency_blue",
            "payroll_frequency_whitecollar": "payroll_frequency_white",
            "primary_operations": "operations",
            "current_hr_policies_summary": "existing_policies",
            "top_hr_compliance_priorities": "priorities",
        }
        for line in combined.splitlines():
            if ":" not in line:
                continue
            key, value = line.split(":", 1)
            normalized = key.strip().lower().replace(" ", "_").replace("(", "").replace(")", "").replace("/", "_")
            normalized = alias_map.get(normalized, normalized)
            snapshot[normalized] = value.strip()
        return normalize_profile_snapshot_keys(snapshot)
    except Exception as e:
        logger.warning(f"Could not fetch company profile snapshot: {e}")
        return {}

def get_company_profile_snapshot_cached(company_id):
    if not company_id:
        return {}
    cache_key = f"profile_snapshot_cache::{company_id}"
    if cache_key in st.session_state:
        return st.session_state.get(cache_key) or {}
    snapshot = get_company_profile_snapshot(company_id)
    st.session_state[cache_key] = snapshot or {}
    return st.session_state[cache_key]

def invalidate_company_profile_snapshot_cache(company_id):
    if not company_id:
        return
    cache_key = f"profile_snapshot_cache::{company_id}"
    if cache_key in st.session_state:
        del st.session_state[cache_key]

def format_profile_memory(profile_snapshot):
    if not profile_snapshot:
        return ""
    lines = [
        "=== COMPANY MEMORY (USE THIS CONTEXT FOR COMPANY-SPECIFIC ANSWERS) ===",
        f"Company Name: {profile_snapshot.get('company_name', '')}",
        f"Sector: {profile_snapshot.get('sector', '')}",
        f"Joint Committees (PC): {profile_snapshot.get('joint_committees', '')}",
        f"Operations: {profile_snapshot.get('operations', '')}",
        f"Headquarters: {profile_snapshot.get('headquarters', '')}",
        f"Countries: {profile_snapshot.get('countries', '')}",
        f"Employees Total: {profile_snapshot.get('employees_total', '')}",
        f"Employees Belgium: {profile_snapshot.get('employees_belgium', '')}",
        f"Contract Types: {profile_snapshot.get('contract_types', '')}",
        f"Weekly Hours: {profile_snapshot.get('weekly_hours', '')}",
        f"Payroll Blue-collar: {profile_snapshot.get('payroll_frequency_blue', '')}",
        f"Payroll White-collar: {profile_snapshot.get('payroll_frequency_white', '')}",
        f"Shift Work: {profile_snapshot.get('shift_work', '')}",
        f"Remote Policy: {profile_snapshot.get('remote_policy', '')}",
        f"Union Presence: {profile_snapshot.get('union_presence', '')}",
        f"Policy Summary: {profile_snapshot.get('existing_policies', '')}",
        f"Compliance Priorities: {profile_snapshot.get('priorities', '')}",
    ]
    return "\n".join(lines)

def compute_profile_completion(profile_snapshot):
    if not profile_snapshot:
        return 0.0
    completed_flag = str(profile_snapshot.get("_profile_completed", "")).strip().lower()
    if completed_flag in {"true", "1", "yes"}:
        return 1.0
    completed = 0
    for field in REQUIRED_PROFILE_FIELDS:
        value = profile_snapshot.get(field, "")
        if isinstance(value, str) and value.strip():
            completed += 1
        elif isinstance(value, (int, float)) and value > 0:
            completed += 1
    return completed / len(REQUIRED_PROFILE_FIELDS)

def index_company_profile(company_id, profile):
    """Store structured company profile as an indexed internal document."""
    try:
        if check_if_document_exists(COMPANY_PROFILE_FILENAME, company_id):
            delete_document(COMPANY_PROFILE_FILENAME, company_id)

        profile_lines = [
            "# Company Profile",
            f"Company Name: {profile.get('company_name', '')}",
            f"Sector: {profile.get('sector', '')}",
            f"Joint Committees (PC): {profile.get('joint_committees', '')}",
            f"Primary Operations: {profile.get('operations', '')}",
            f"Industry Cluster: {profile.get('industry_cluster', '')}",
            f"Headquarters: {profile.get('headquarters', '')}",
            f"Countries of Operation: {profile.get('countries', '')}",
            f"Preferred Language: {profile.get('language', '')}",
            "",
            "## Workforce",
            f"Total Employees: {profile.get('employees_total', '')}",
            f"Employees in Belgium: {profile.get('employees_belgium', '')}",
            f"Contract Types: {profile.get('contract_types', '')}",
            f"Union Delegation Present: {profile.get('union_presence', '')}",
            "",
            "## Working Conditions",
            f"Standard Weekly Hours: {profile.get('weekly_hours', '')}",
            f"Shift/Night/Weekend Work: {profile.get('shift_work', '')}",
            f"Remote Work Policy: {profile.get('remote_policy', '')}",
            f"Payroll Frequency (Blue-collar): {profile.get('payroll_frequency_blue', '')}",
            f"Payroll Frequency (White-collar): {profile.get('payroll_frequency_white', '')}",
            "",
            "## Policies & Context",
            f"Current HR Policies Summary: {profile.get('existing_policies', '')}",
            f"Top HR Compliance Priorities: {profile.get('priorities', '')}",
            f"Open Legal Questions: {profile.get('open_questions', '')}",
            "",
            "## Industry Specific",
            f"Industry Specific Details: {profile.get('industry_specific_details', '')}",
        ]
        profile_text = "\n".join(profile_lines).strip()
        chunks = recursive_chunk_text(profile_text, chunk_size=700, overlap=100)
        if not chunks:
            return False

        stored = 0
        for i in range(0, len(chunks), 20):
            batch = chunks[i:i+20]
            vectors = get_embeddings_batch(batch)
            if not vectors:
                continue
            payload = []
            for j, vec in enumerate(vectors):
                if isinstance(vec, list) and len(vec) > 300:
                    payload.append({
                        "content": batch[j],
                        "metadata": {
                            "company_id": company_id,
                            "filename": COMPANY_PROFILE_FILENAME,
                            "is_active": True,
                            "title": "Company Profile",
                            "author": profile.get("company_name", "Company"),
                            "document_type": "company_profile",
                        },
                        "embedding": vec
                    })
            if payload:
                supabase.table("document_chunks").insert(payload).execute()
                stored += len(payload)

        if stored == 0:
            return False

        # Store a compact JSON snapshot for reliable progress calculation and hydration
        try:
            # Remove old profile snapshots for this company to avoid stale progress reads
            supabase.table("document_chunks").delete().eq(
                "metadata->>company_id", company_id
            ).eq(
                "metadata->>document_type", COMPANY_PROFILE_SNAPSHOT_TYPE
            ).execute()

            snapshot_json = json.dumps(profile, ensure_ascii=True)
            snapshot_text = f"__PROFILE_JSON__:{snapshot_json}"
            snapshot_vecs = get_embeddings_batch([snapshot_text])
            if snapshot_vecs and isinstance(snapshot_vecs[0], list) and len(snapshot_vecs[0]) > 300:
                supabase.table("document_chunks").insert({
                    "content": snapshot_text,
                    "metadata": {
                        "company_id": company_id,
                        "filename": COMPANY_PROFILE_FILENAME,
                        "is_active": True,
                        "title": "Company Profile Snapshot",
                        "document_type": COMPANY_PROFILE_SNAPSHOT_TYPE,
                    },
                    "embedding": snapshot_vecs[0],
                }).execute()
        except Exception as e:
            logger.warning(f"Could not store compact profile snapshot: {e}")

        registered = register_document(COMPANY_PROFILE_FILENAME, company_id, {"title": "Company Profile"})
        if registered:
            invalidate_company_profile_snapshot_cache(company_id)
        return bool(registered)
    except Exception as e:
        logger.error(f"Failed to index company profile: {e}")
        return False

def process_and_store_legal_document(file, legal_type, source_code="", topic="", effective_date=""):
    """Index legal sources into legal_knowledge for cross-company legal retrieval."""
    if not file or not hasattr(file, 'name') or not hasattr(file, 'size'):
        return "error"
    if file.size > 50 * 1024 * 1024:
        return "error"

    clean_name = sanitize_filename(file.name)
    try:
        text, metadata, ext = extract_text_and_metadata(file)
        if ext not in {"pdf", "docx", "xlsx", "pptx"}:
            return "unsupported"
    except Exception as e:
        logger.error(f"Legal extraction failed for {clean_name}: {e}")
        return "error"

    if not text or len(text.strip()) < 10:
        return "empty"

    chunks = recursive_chunk_text(text)
    if not chunks:
        return "empty"

    normalized_source = re.sub(r"[^A-Za-z0-9_-]", "_", (source_code or "").strip().upper())
    source = "BELGIAN_LAW" if legal_type == "law" else (normalized_source or "CAO_UNKNOWN")
    category = "legal_foundation" if legal_type == "law" else "sector_agreement"
    legal_tier = 1 if legal_type == "law" else 2

    stored = 0
    for i in range(0, len(chunks), 20):
        batch = chunks[i:i+20]
        vectors = get_embeddings_batch(batch)
        if not vectors:
            continue
        payload = []
        for j, vec in enumerate(vectors):
            if isinstance(vec, list) and len(vec) > 300:
                chunk_text = batch[j]
                hash_input = f"{source}|{clean_name}|{i + j}|{chunk_text}".encode("utf-8", errors="ignore")
                payload.append({
                    "content": chunk_text,
                    "summary": chunk_text[:500],
                    "content_hash": hashlib.md5(hash_input).hexdigest(),
                    "metadata": {
                        "source": source,
                        "category": category,
                        "topic": topic or metadata.get("title") or clean_name,
                        "effective_date": effective_date or None,
                        "uploaded_filename": clean_name,
                        "legal_tier": legal_tier,
                    },
                    "embedding": vec,
                })
        if payload:
            try:
                supabase.table("legal_knowledge").upsert(payload, on_conflict="content_hash").execute()
                stored += len(payload)
            except Exception as e:
                logger.error(f"Failed storing legal batch for {clean_name}: {e}")
    return "success" if stored > 0 else "error"

def process_and_store_document(file, company_id, force_overwrite=False):
    """Process and store document with support for PDF, DOCX, XLSX, PPTX."""
    # Validate file object
    if not file or not hasattr(file, 'name') or not hasattr(file, 'size'):
        logger.error("Invalid file object provided")
        return "error"
    
    # Check file size (50MB limit)
    MAX_FILE_SIZE = 50 * 1024 * 1024  # 50MB
    if file.size > MAX_FILE_SIZE:
        logger.error(f"File too large: {file.size} bytes (max: {MAX_FILE_SIZE} bytes)")
        return "error"
    
    clean_name = sanitize_filename(file.name)
    logger.info(f"Starting to process document: {clean_name}")
    
    if check_if_document_exists(clean_name, company_id):
        if not force_overwrite: 
            logger.info(f"Document already exists: {clean_name}")
            return "exists"
        else: 
            logger.info(f"Overwriting existing document: {clean_name}")
            delete_document(clean_name, company_id)

    text = ""
    file_metadata = {"title": clean_name}
    try:
        text, file_metadata, ext = extract_text_and_metadata(file)
        logger.info(f"Processing file type: {ext}")
        if ext not in {"pdf", "docx", "xlsx", "pptx"}:
            logger.warning(f"Unsupported file type: {ext}")
            return "unsupported"
    except Exception as e:
        logger.error(f"Extraction failed for {clean_name}: {e}")
        return "error"
    
    if not text or len(text.strip()) < 10: 
        logger.warning(f"Document empty or too short: {clean_name}")
        return "empty"
    
    logger.info(f"Extracted {len(text)} characters from {clean_name}")
    
    # === VISUAL RAG: Extract and describe images ===
    try:
        from services.vision_service import get_visual_context
        file.seek(0)
        # Process top 35 images to capture medium-sized infographics (hairnets, icons)
        # 35 images handles documents with many small icons while respecting rate limits
        visual_context = get_visual_context(file, FIXED_GROQ_KEY, max_images=35)
        if visual_context:
            text = text + visual_context
            logger.info(f"Added visual context: {len(visual_context)} chars")
    except Exception as e:
        logger.warning(f"Visual context extraction failed (non-critical): {e}")
    # === END VISUAL RAG ===

    try:
        file.seek(0)
        # Sanitize company_id for storage path (remove emoji and special chars)
        safe_company_id = re.sub(r'[^\w\-]', '_', company_id)
        # Determine correct Content-Type to prevent browser rendering as text
        mime_type = get_mime_type_for_extension(ext)
        
        supabase.storage.from_("documents").upload(
            f"{safe_company_id}/{clean_name}", 
            file.read(), 
            {"upsert": "true", "contentType": mime_type}
        )
    except Exception as e:
        logger.warning(f"Storage upload failed (non-critical): {e}")

    chunks = recursive_chunk_text(text)
    logger.info(f"Created {len(chunks)} chunks from {clean_name}")
    
    chunks_stored = 0
    failed_batches = []
    
    for i in range(0, len(chunks), 20):
        batch = chunks[i:i+20]
        vectors = get_embeddings_batch(batch)
        
        if not vectors:
            logger.warning(f"Embedding generation returned None for batch starting at {i}")
            failed_batches.append({"start": i, "reason": "embedding_failed"})
            continue
            
        payload = []
        for j, vec in enumerate(vectors):
            if isinstance(vec, list) and len(vec) > 300:
                payload.append({
                    "content": batch[j],
                    "metadata": {
                        "company_id": company_id, "filename": clean_name, "is_active": True,
                        "title": file_metadata.get("title"), "author": file_metadata.get("author")
                    },
                    "embedding": vec
                })
        
        if payload:
            # Retry with exponential backoff
            max_retries = 3
            for attempt in range(max_retries):
                try:
                    supabase.table("document_chunks").insert(payload).execute()
                    chunks_stored += len(payload)
                    break  # Success!
                except Exception as e:
                    if attempt == max_retries - 1:
                        logger.error(f"FAILED after {max_retries} attempts for batch {i}: {e}")
                        failed_batches.append({"start": i, "reason": str(e)})
                    else:
                        wait_time = 2 ** (attempt + 1)  # 2s, 4s, 8s
                        logger.warning(f"Retry {attempt + 1}/{max_retries} for batch {i}, waiting {wait_time}s...")
                        time.sleep(wait_time)
    
    if failed_batches:
        logger.warning(f"Document {clean_name}: {len(failed_batches)} batch(es) failed: {failed_batches}")
    
    logger.info(f"Stored {chunks_stored} chunks for {clean_name}")
    
    if register_document(clean_name, company_id, file_metadata):
        logger.info(f"Successfully registered document: {clean_name}")
        return "success"
    return "error"

# --- [FIX 2] INTELLIGENT SEARCH EXPANSION ---
def get_relevant_context(query, company_id):
    normalized_query = normalize_query(query)
    search_query = normalized_query
    
    try:
        # EXPANSION PROMPT: Explicitly ask for split and joined variations
        # This fixes "koffie machine" vs "koffiemachine" vs "coffee machine"
        expansion_prompt = (
            f"You are a search optimizer. User Query: '{normalized_query}'. "
            f"Generate a boolean search string that includes: "
            f"1. English, Dutch, and French translations. "
            f"2. CRITICAL: For every compound word, output BOTH the joined version (e.g. 'koffiemachine') AND the split version (e.g. 'coffee machine'). "
            f"Output ONLY the terms separated by ' OR '."
        )
        
        expansion_models = ["llama-3.3-70b-versatile", "llama-3.1-8b-instant"]
        for model in expansion_models:
            resp = requests.post(
                "https://api.groq.com/openai/v1/chat/completions",
                headers={"Authorization": f"Bearer {FIXED_GROQ_KEY}"},
                json={"model": model, "messages": [{"role": "user", "content": expansion_prompt}], "temperature": 0.1, "max_tokens": 100},
                timeout=10
            )
            if resp.status_code == 200:
                search_query = resp.json()['choices'][0]['message']['content']
                break
            elif resp.status_code == 429: continue
            else: break
    except: pass

    # Get Vector Embedding for the ORIGINAL normalized query (to capture intent)
    vectors = get_embeddings_batch([normalized_query])
    if not vectors: return "", []

    try:
        # Use the EXPANDED query for Full Text Search (FTS) to catch the English keywords
        params = {"query_embedding": vectors[0], "match_threshold": 0.15, "match_count": 15, "filter_company_id": company_id, "query_text": search_query}
        res = supabase.rpc("match_documents_hybrid", params).execute()
        context_str = ""
        sources = []
        for m in res.data:
            context_str += f"-- SOURCE: {m['metadata']['filename']} --\n{m['content']}\n\n"
            if m['metadata']['filename'] not in sources: sources.append(m['metadata']['filename'])
        return context_str, sources
    except: return "", []

def ask_groq(context, history, query):
    """Ask Groq LLM with multilingual support and smart source citation."""
    system_prompt = """You are FRIDAY, an expert multilingual HR assistant.

CRITICAL RULES:
1. LANGUAGE: Always respond in the SAME LANGUAGE as the user's query.
   - Dutch query → Dutch response
   - English query → English response
   - French query → French response

2. SOURCES: Answer based ONLY on the provided CONTEXT.
   - If the answer requires info from multiple documents, combine them intelligently
   - ONLY cite documents you actually used in your answer
   - Cite using the exact filename (e.g., "according to onboarding_guide.pdf")
   - Do NOT cite documents that don't contain relevant info

3. HONESTY: If the CONTEXT doesn't contain the answer, say so clearly.

4. FORMAT: Use Markdown for tables and structure when helpful.

5. FORMS & DOWNLOADS:
   - If the context contains a "RECOMMENDED FORMS" section with links, you MUST mention them.
   - Example: "You can download the form here: [link]."
   - Only recommend forms that are explicitly listed in the context.

6. LEGAL CITATIONS (MANDATORY):
   - The CONTEXT section contains Belgian labor law articles. Each article has its number and law name in the header.
   - When using ANY information from a law article, you MUST cite the EXACT article number and law name from that article's header.
   - Format: "Conform [artikelnummer] van de [wetnaam] ([datum]): ..."
   - Extract the article number FROM THE CONTEXT HEADER, do NOT make up article numbers.
   - If the law text directly answers the question, QUOTE the relevant passage.
   - When multiple articles are relevant, cite ALL of them with their correct numbers.
   - NEVER cite an article number that is not present in the provided context."""
    
    messages = [{"role": "system", "content": system_prompt}]
    for msg in history[-4:]: messages.append({"role": msg["role"], "content": msg["content"]})
    if context: messages.append({"role": "user", "content": f"CONTEXT (Filenames in headers):\n{context}"})
    messages.append({"role": "user", "content": query})

    models = ["llama-3.3-70b-versatile", "llama-3.1-8b-instant"]
    for model in models:
        try:
            resp = requests.post(
                "https://api.groq.com/openai/v1/chat/completions",
                headers={"Authorization": f"Bearer {FIXED_GROQ_KEY}"},
                json={"model": model, "messages": messages, "temperature": 0.1},
                timeout=30
            )
            if resp.status_code == 200: return resp.json()['choices'][0]['message']['content']
            elif resp.status_code == 429: continue
        except: pass
    return "⚠️ Service currently limited. Please try again."

# --- CHAT HISTORY & PERSISTENCE ---

def load_chat_history(chat_id):
    try: return supabase.table("messages").select("*").eq("chat_id", chat_id).order("created_at").execute().data
    except: return []

def save_message(chat_id, role, content, company_id, sources=None):
    try:
        import json
        
        # Ensure sources is properly serialized
        sources_data = sources if sources is not None else {}
        
        # CRITICAL FIX: For database storage, merge the new 3-tier structure back into 
        # the legacy 2-key format (legal_sources, company_sources) for backward compatibility
        # The database might have constraints expecting only these two keys
        if isinstance(sources_data, dict):
            # If sources has the new structure with law/sector/company tiers
            if "law_sources" in sources_data or "sector_sources" in sources_data:
                # Deduplicate: legal_sources already contains law+sector+company_legal
                # from rag_controller, so use it directly instead of re-concatenating
                all_legal = (
                    sources_data.get("law_sources", []) + 
                    sources_data.get("sector_sources", []) +
                    sources_data.get("legal_sources", [])
                )
                # Remove duplicates while preserving order
                seen = set()
                deduped_legal = []
                for s in all_legal:
                    if s not in seen:
                        seen.add(s)
                        deduped_legal.append(s)
                legacy_sources = {
                    "legal_sources": deduped_legal,
                    "company_sources": sources_data.get("company_sources", [])
                }
                sources_data = legacy_sources
                logger.info(f"Converted 3-tier sources to legacy format for DB storage")
        
        # Validate that sources is JSON-serializable
        try:
            json.dumps(sources_data)
        except (TypeError, ValueError) as e:
            logger.error(f"Sources not JSON-serializable: {e}, sources={sources_data}")
            sources_data = {}
        
        # CRITICAL FIX: The DB `sources` column is typed as a JSON array.
        # Wrap the dict in a list so Postgres accepts it as a valid JSON array.
        # Format: [{"legal_sources": [...], "company_sources": [...]}]
        if isinstance(sources_data, dict):
            sources_for_db = [sources_data]
        elif isinstance(sources_data, list):
            sources_for_db = sources_data
        else:
            sources_for_db = []
        
        supabase.table("messages").insert({
            "chat_id": chat_id, 
            "role": role, 
            "content": content, 
            "sources": sources_for_db, 
            "company_id": company_id 
        }).execute()
        logger.info(f"Successfully saved {role} message to chat {chat_id}")
    except Exception as e:
        logger.error(f"Failed to save message: {e}, chat_id={chat_id}, role={role}, has_sources={sources is not None}")
        # Don't re-raise to prevent chat from breaking
        pass

def get_recent_chats(company_id):
    try:
        res = supabase.table("messages").select("chat_id, content, created_at").eq("company_id", company_id).order("created_at", desc=True).limit(50).execute()
        seen_ids = set()
        unique_chats = []
        for msg in res.data:
            if msg['chat_id'] not in seen_ids:
                seen_ids.add(msg['chat_id'])
                unique_chats.append({"id": msg['chat_id'], "title": msg['content'][:30] + "..."})
        return unique_chats[:10]
    except: return []

def delete_chat(chat_id, company_id):
    try:
        supabase.table("messages").delete().eq("chat_id", chat_id).eq("company_id", company_id).execute()
        return True
    except: return False

def get_dynamic_greeting():
    hour = datetime.datetime.now().hour
    if 5 <= hour < 12: return "Good morning."
    elif 12 <= hour < 17: return "Good afternoon."
    elif 17 <= hour < 21: return "Good evening."
    else: return "Hello."

# --- SOURCE DISPLAY HELPER ---
def render_sources_html(sources) -> str:
    """Render sources with 3-tier legal hierarchy distinction.
    
    Handles all formats:
    - dict with law_sources, sector_sources, company_sources (live display)
    - dict with legal_sources and company_sources (legacy or DB round-trip)
    - list wrapping a dict [{...}] (DB storage format)
    - flat list of strings (very old messages)
    """
    # DB format: sources stored as ['{"legal_sources":[...],"company_sources":[...]}']
    # Can be either a JSON string or a dict inside the list — handle both
    if isinstance(sources, list):
        if len(sources) == 1:
            item = sources[0]
            # Case A: [{dict}] — unwrap directly
            if isinstance(item, dict):
                sources = item
            # Case B: ['{"legal_sources":...}'] — JSON string, parse it
            elif isinstance(item, str) and item.strip().startswith('{'):
                import json
                try:
                    sources = json.loads(item)
                except (json.JSONDecodeError, ValueError):
                    return _render_flat_sources(sources)
            else:
                return _render_flat_sources(sources)
        elif all(isinstance(s, str) for s in sources):
            return _render_flat_sources(sources)  # Truly flat legacy list
        else:
            return _render_flat_sources([str(s) for s in sources])
    
    if not isinstance(sources, dict):
        return _render_empty_sources()
    
    # Extract tiers — support both new (law/sector) and legacy (legal_sources) keys
    law = sources.get("law_sources", [])
    sector = sources.get("sector_sources", [])
    company = sources.get("company_sources", [])
    
    # Backward compat: if no law/sector keys, fall back to legal_sources
    if not law and not sector and sources.get("legal_sources"):
        law = sources.get("legal_sources", [])
    
    total = len(law) + len(sector) + len(company)
    
    if total == 0:
        return _render_empty_sources()
    
    # Build sections by tier — IMPORTANT: no leading whitespace!
    # Streamlit's markdown parser treats 4+ spaces as code blocks
    law_html = ""
    if law:
        law_pills = ''.join([f'<span class="source-pill legal-source-pill">⚖️ {s}</span>' for s in law])
        law_html = f'<div class="source-group-header legal">⚖️ Federal Law</div><div class="sources-list">{law_pills}</div>'
    
    sector_html = ""
    if sector:
        sector_pills = ''.join([f'<span class="source-pill sector-source-pill">📋 {s}</span>' for s in sector])
        sector_html = f'<div class="source-group-header sector">📋 Sector Agreement</div><div class="sources-list">{sector_pills}</div>'
    
    company_html = ""
    if company:
        company_pills = ''.join([f'<span class="source-pill company-source-pill">📄 {s}</span>' for s in company])
        company_html = f'<div class="source-group-header company">📄 Company Policy</div><div class="sources-list">{company_pills}</div>'
    
    if total == 1:
        return f'<div class="sources-section">{law_html}{sector_html}{company_html}</div>'
    
    return f'<div class="sources-section"><details><summary>📚 {total} sources used</summary>{law_html}{sector_html}{company_html}</details></div>'


def _render_empty_sources() -> str:
    """Render the 'no sources' placeholder."""
    return '<div class="sources-section"><div class="no-sources"><span>💭</span><span>This response was generated without referencing any documents</span></div></div>'


def _render_flat_sources(sources: list) -> str:
    """Legacy renderer for old messages stored as flat list of strings."""
    if not sources:
        return _render_empty_sources()
    pills = ''.join([f'<span class="source-pill company-source-pill">📄 {s}</span>' for s in sources])
    if len(sources) == 1:
        return f'<div class="sources-section"><div class="sources-list">{pills}</div></div>'
    return f'<div class="sources-section"><details><summary>📚 {len(sources)} sources used</summary><div class="sources-list">{pills}</div></details></div>'


# --- UI PAGES ---
def render_sidebar():
    with st.sidebar:
        # Logo
        st.markdown('<div style="display: flex; justify-content: center; align-items: center; height: 100px; margin-bottom: 28px;"><span class="logo-animated" style="font-family: \'Playfair Display\', serif; font-size: 56px; font-weight: 800; color: #1A3C34; letter-spacing: -1.2px;">Friday</span></div>', unsafe_allow_html=True)

        profile_snapshot = get_company_profile_snapshot_cached(st.session_state.company_id)
        completion_ratio = compute_profile_completion(profile_snapshot)
        st.markdown("### Company Setup")
        st.progress(completion_ratio, text=f"{int(completion_ratio * 100)}% complete")
        if completion_ratio < 1.0:
            missing = []
            for f in REQUIRED_PROFILE_FIELDS:
                v = profile_snapshot.get(f, "")
                if not (isinstance(v, str) and v.strip()) and not (isinstance(v, (int, float)) and v > 0):
                    missing.append(f)
            with st.expander("Profile completeness details"):
                st.caption("Missing fields")
                st.code(", ".join(missing) if missing else "None")
        if st.button("✎ Update Company Profile", use_container_width=True, type="secondary"):
            st.session_state.onboarding_required = True
            st.session_state.onboarding_step = 1
            st.rerun()
        st.markdown("---")

        # Navigation
        st.markdown("### Menu")
        # [FIX 1] Button selection logic is correct here, CSS handles the color
        if st.button("◉  Chat", use_container_width=True, type="primary" if st.session_state.view == "chat" else "secondary"):
            st.session_state.view = "chat"; st.rerun()
        if st.button("◎  Documents", use_container_width=True, type="primary" if st.session_state.view == "documents" else "secondary"):
            st.session_state.view = "documents"; st.rerun()

        st.markdown("---")

        if st.session_state.view == "chat":
            if st.button("＋  New Chat", use_container_width=True, type="secondary"): 
                create_new_chat(); st.rerun()

            st.markdown("### Recent Chats")
            recent = get_recent_chats(st.session_state.company_id)
            if not recent:
                st.caption("No conversations yet")
            else:
                for chat in recent:
                    is_active = st.session_state.current_chat_id == chat['id']
                    with st.container():
                        col1, col2 = st.columns([6, 1])
                        with col1:
                            # Dynamic button type ensures CSS triggers White or Black text
                            btn_type = "primary" if is_active else "secondary"
                            if st.button(f"{chat['title']}", key=f"chat_{chat['id']}", use_container_width=True, type=btn_type):
                                st.session_state.current_chat_id = chat['id']
                                st.rerun()
                        with col2:
                            if st.button("×", key=f"del_{chat['id']}"):
                                delete_chat(chat['id'], st.session_state.company_id)
                                if st.session_state.current_chat_id == chat['id']: create_new_chat()
                                st.rerun()

        st.markdown("---")
        st.markdown('<div class="logout-btn">', unsafe_allow_html=True)
        if st.button("↪  Sign Out", use_container_width=True): 
            st.session_state.clear(); st.rerun()
        st.markdown('</div>', unsafe_allow_html=True)

def create_new_chat(): st.session_state.current_chat_id = str(uuid.uuid4())

def handle_query(query):
    save_message(st.session_state.current_chat_id, "user", query, st.session_state.company_id)
    with st.chat_message("user", avatar="👤"): st.write(query)
    
    with st.chat_message("assistant", avatar="⚡"):
        msg_placeholder = st.empty()
        # Custom thinking spinner
        spinner_html = '''
        <div class="thinking-spinner">
            <span></span><span></span><span></span>
            <span class="thinking-text">Friday is thinking...</span>
        </div>
        '''
        msg_placeholder.markdown(spinner_html, unsafe_allow_html=True)
        
        history = load_chat_history(st.session_state.current_chat_id)
        
        # Elite RAG Pipeline
        import asyncio
        context, source_dict = asyncio.run(get_context_with_strategy(
            raw_query=query,
            company_id=st.session_state.company_id,
            supabase=supabase,
            groq_api_key=FIXED_GROQ_KEY,
            get_embeddings_fn=get_embeddings_batch,
            hf_api_key=HF_API_KEY,
            top_k=10
        ))

        # Always include current company profile memory in the final context.
        profile_snapshot = get_company_profile_snapshot_cached(st.session_state.company_id)
        profile_memory = format_profile_memory(profile_snapshot)
        if profile_memory:
            context = f"{profile_memory}\n\n{context}" if context else profile_memory
        
        response = ask_groq(context, history, query)
        
        # source_dict = {"legal_sources": [...], "company_sources": [...]}
        display_sources = source_dict if isinstance(source_dict, dict) else {"legal_sources": [], "company_sources": []}
        if profile_memory and isinstance(display_sources, dict):
            cs = display_sources.get("company_sources", [])
            if "Company Profile Memory" not in cs:
                cs = ["Company Profile Memory"] + cs
                display_sources["company_sources"] = cs
        save_message(st.session_state.current_chat_id, "assistant", response, st.session_state.company_id, display_sources)
        msg_placeholder.empty()
        st.markdown(response)
        
        # Dual-source display
        st.markdown(render_sources_html(display_sources), unsafe_allow_html=True)
        
        # Confidence scoring (math-based, no extra API calls)
        from services.agentic.confidence import calculate_confidence, format_confidence_html
        total_sources = len(display_sources.get("legal_sources", [])) + len(display_sources.get("company_sources", []))
        confidence = calculate_confidence(
            sources_count=total_sources,
            context_length=len(context) if context else 0,
            response_length=len(response) if response else 0
        )
        st.markdown(format_confidence_html(confidence), unsafe_allow_html=True)
        
    st.rerun()

def chat_page():
    if not st.session_state.current_chat_id: create_new_chat()
    history = load_chat_history(st.session_state.current_chat_id)

    if st.session_state.get("show_welcome_anim"):
        st.markdown(
            '<div class="welcome-shell"><div class="welcome-typewriter">Welcome to Friday</div></div>',
            unsafe_allow_html=True
        )
        st.session_state.show_welcome_anim = False
        time.sleep(1.1)
        st.rerun()

    if not history:
        greeting = get_dynamic_greeting()
        st.markdown(f"""
        <div style="text-align: center; margin-top: 80px; margin-bottom: 48px;">
            <h1 class="greeting-title">{greeting}</h1>
            <p class="greeting-subtitle">How can FRIDAY help you with HR tasks today?</p>
        </div>
        """, unsafe_allow_html=True)

    for msg in history:
        with st.chat_message(msg["role"], avatar="⚡" if msg["role"] == "assistant" else "👤"):
            st.write(msg["content"])
            if msg["role"] == "assistant":
                sources = msg.get("sources", []) or []
                st.markdown(render_sources_html(sources), unsafe_allow_html=True)

    if prompt := st.chat_input("Ask FRIDAY anything..."):
        handle_query(prompt)

def documents_page():
    st.markdown("""
    <div style="margin-bottom: 32px;">
        <h1 style="font-family: 'Playfair Display', serif; font-size: 48px; font-weight: 700; color: #1A3C34;">Knowledge Base</h1>
    </div>
    """, unsafe_allow_html=True)
    company_tab, legal_tab = st.tabs(["Company Documents", "Belgian Labor Law & CAO"])

    with company_tab:
        col1, col2 = st.columns([1, 1], gap="large")
        with col1:
            st.markdown("### Upload Internal Documents")
            uploaded_files = st.file_uploader(
                "PDF, Word, Excel, PPTX",
                type=["pdf", "docx", "xlsx", "pptx"],
                accept_multiple_files=True,
                key="company_docs_uploader"
            )

            c_check, c_btn = st.columns([1, 1])
            with c_check:
                force_overwrite = st.checkbox("Overwrite existing?", key="company_overwrite")
            with c_btn:
                if uploaded_files and st.button("Start Indexing", type="primary", use_container_width=True, key="company_index_btn"):
                    results = {"success": 0, "exists": 0, "error": 0}
                    progress_bar = st.progress(0)
                    status = st.empty()

                    for idx, f in enumerate(uploaded_files):
                        status.text(f"Indexing {f.name}...")
                        res = process_and_store_document(f, st.session_state.company_id, force_overwrite)
                        if res == "success":
                            results["success"] += 1
                        elif res == "exists":
                            results["exists"] += 1
                        else:
                            results["error"] += 1
                        progress_bar.progress((idx + 1) / len(uploaded_files))

                    status.empty()
                    progress_bar.empty()
                    if results["success"] > 0:
                        st.success(f"✅ {results['success']} files indexed successfully.")
                    if results["exists"] > 0:
                        st.info(f"ℹ️ {results['exists']} files already existed.")
                    if results["error"] > 0:
                        st.error(f"❌ {results['error']} files failed to process.")
                    time.sleep(1)
                    st.rerun()

        with col2:
            docs = get_all_documents(st.session_state.company_id)
            chunk_counts = get_chunk_counts(st.session_state.company_id)
            missing_chunks = sum(1 for doc in docs if chunk_counts.get(doc['filename'], 0) == 0)

            st.markdown(f"### Indexed Files ({len(docs)})")
            if missing_chunks > 0:
                st.warning(f"⚠️ {missing_chunks} document(s) have no searchable chunks. Check 'Overwrite existing?' and re-upload them.")

            if not docs:
                st.info("No documents indexed yet.")
            else:
                for doc in docs:
                    chunks = chunk_counts.get(doc['filename'], 0)
                    status_icon = "🟢" if doc['is_active'] else "⚪"
                    chunk_warning = " ⚠️" if chunks == 0 else ""
                    file_ext = doc['filename'].split('.')[-1].upper()
                    display_name = doc['filename'][:37] + "..." if len(doc['filename']) > 40 else doc['filename']

                    col_status, col_name, col_chunks, col_toggle, col_delete = st.columns([0.5, 3.5, 0.8, 0.5, 0.5])
                    with col_status:
                        st.markdown(f"<div style='padding-top: 8px;'>{status_icon}</div>", unsafe_allow_html=True)
                    with col_name:
                        st.markdown(f'''
                        <div style="display: flex; align-items: center; gap: 8px; padding: 8px 0;">
                            <span style="font-weight: 500; color: #1A3C34;">{file_ext}</span>
                            <span title="{doc['filename']}" style="overflow: hidden; text-overflow: ellipsis; white-space: nowrap;">{display_name}</span>
                        </div>
                        ''', unsafe_allow_html=True)
                    with col_chunks:
                        chunk_color = "#e74c3c" if chunks == 0 else "#1A3C34"
                        st.markdown(f'''
                        <div style="padding: 8px 0; font-size: 12px; color: {chunk_color};" title="Number of searchable chunks">
                            {chunks} chunks{chunk_warning}
                        </div>
                        ''', unsafe_allow_html=True)
                    with col_toggle:
                        if st.button("⏸" if doc['is_active'] else "▶", key=f"pause_{doc['id']}", help="Pause/Resume"):
                            toggle_document_status(doc['filename'], st.session_state.company_id, doc['is_active'])
                            st.rerun()
                    with col_delete:
                        if st.button("🗑", key=f"del_doc_{doc['id']}", help="Delete"):
                            delete_document(doc['filename'], st.session_state.company_id)
                            st.rerun()
                    st.markdown("<hr style='margin: 4px 0; border: none; border-top: 1px solid #eee;'>", unsafe_allow_html=True)

    with legal_tab:
        st.markdown('<div class="legal-upload-card">', unsafe_allow_html=True)
        st.markdown("### Upload Belgian Labor Law and CAOs")
        st.caption("Use this to quickly add legal references FRIDAY can cite directly in answers.")

        legal_type = st.selectbox(
            "Document type",
            options=[("law", "Belgian Federal Law"), ("cao", "CAO / Sector Agreement")],
            format_func=lambda x: x[1],
            key="legal_type"
        )[0]
        source_code = st.text_input(
            "Source code",
            value="BELGIAN_LAW" if legal_type == "law" else "",
            help="Examples: BELGIAN_LAW, PC200, CAO_XXX",
            key="legal_source_code"
        )
        topic = st.text_input("Topic (optional)", placeholder="Employment contracts, Working time, Wages...", key="legal_topic")
        effective_date = st.text_input("Effective date (YYYY-MM-DD, optional)", key="legal_effective_date")
        legal_files = st.file_uploader(
            "Upload legal files",
            type=["pdf", "docx", "xlsx", "pptx"],
            accept_multiple_files=True,
            key="legal_docs_uploader"
        )

        if legal_files and st.button("Index Legal Sources", type="primary", use_container_width=True, key="legal_index_btn"):
            results = {"success": 0, "error": 0}
            progress_bar = st.progress(0)
            status = st.empty()
            date_value = effective_date.strip()

            for idx, f in enumerate(legal_files):
                status.text(f"Indexing {f.name}...")
                res = process_and_store_legal_document(
                    file=f,
                    legal_type=legal_type,
                    source_code=source_code,
                    topic=topic,
                    effective_date=date_value,
                )
                if res == "success":
                    results["success"] += 1
                else:
                    results["error"] += 1
                progress_bar.progress((idx + 1) / len(legal_files))

            status.empty()
            progress_bar.empty()
            if results["success"] > 0:
                st.success(f"✅ {results['success']} legal file(s) indexed.")
            if results["error"] > 0:
                st.error(f"❌ {results['error']} legal file(s) failed.")

        st.markdown('</div>', unsafe_allow_html=True)

def get_industry_specific_payload(industry_cluster):
    """Return (details_text, is_valid) for the selected industry-specific step."""
    if industry_cluster == "Hospitality":
        weekend_days = st.multiselect("Weekend opening days *", ["Saturday", "Sunday"], key="ob_weekend_days")
        night_service = st.selectbox("Night service after 22:00 *", ["No", "Occasionally", "Regularly"], key="ob_night_service")
        tipped_roles = st.selectbox("Tipped/front-of-house roles *", ["Yes", "No"], key="ob_tipped_roles")
        details = f"Weekend days: {', '.join(weekend_days)}; Night service: {night_service}; Tipped roles: {tipped_roles}"
        return details, len(weekend_days) > 0
    if industry_cluster == "Manufacturing":
        shift_pattern = st.selectbox("Shift model *", ["No shifts", "2 shifts", "3 shifts", "Continuous"], key="ob_shift_pattern")
        hazardous_work = st.selectbox("Hazardous work/materials *", ["No", "Yes"], key="ob_hazardous_work")
        temp_peaks = st.selectbox("Seasonal temp-worker peaks *", ["No", "Yes"], key="ob_temp_peaks")
        return f"Shift model: {shift_pattern}; Hazardous work: {hazardous_work}; Temp peaks: {temp_peaks}", True
    if industry_cluster == "Retail":
        store_count = st.number_input("Number of stores/sites *", min_value=1, step=1, key="ob_store_count")
        sunday_open = st.selectbox("Sunday opening *", ["No", "Yes"], key="ob_sunday_open")
        late_openings = st.selectbox("Late openings after 20:00 *", ["No", "Occasionally", "Regularly"], key="ob_late_openings")
        return f"Stores: {int(store_count)}; Sunday opening: {sunday_open}; Late openings: {late_openings}", True
    if industry_cluster == "Logistics/Transport":
        drivers_count = st.number_input("Number of drivers *", min_value=0, step=1, key="ob_drivers_count")
        cross_border = st.selectbox("Cross-border operations *", ["No", "Yes"], key="ob_cross_border")
        warehouse_24_7 = st.selectbox("24/7 warehouse operations *", ["No", "Yes"], key="ob_warehouse_24_7")
        return f"Drivers: {int(drivers_count)}; Cross-border: {cross_border}; 24/7 warehouse: {warehouse_24_7}", True
    if industry_cluster == "Healthcare":
        on_call = st.selectbox("On-call duty *", ["No", "Yes"], key="ob_on_call")
        weekend_care = st.selectbox("Weekend care staffing *", ["No", "Yes"], key="ob_weekend_care")
        regulated_titles = st.selectbox("Regulated medical roles *", ["No", "Yes"], key="ob_regulated_titles")
        return f"On-call: {on_call}; Weekend care: {weekend_care}; Regulated roles: {regulated_titles}", True
    if industry_cluster == "Professional Services":
        billable_hours = st.selectbox("Billable-hours model *", ["No", "Yes"], key="ob_billable_hours")
        client_site_work = st.selectbox("Client-site work *", ["No", "Occasionally", "Regularly"], key="ob_client_site_work")
        overtime_policy = st.selectbox("Overtime compensation policy defined *", ["No", "Yes"], key="ob_overtime_policy")
        return f"Billable model: {billable_hours}; Client-site work: {client_site_work}; Overtime policy: {overtime_policy}", True
    details = st.text_area(
        "Industry-specific constraints *",
        placeholder="Describe schedules, regulations, or collective requirements.",
        height=90,
        key="ob_industry_other_details"
    )
    return details, bool(details.strip())

def validate_onboarding_step(step):
    checks = {
        1: [
            st.session_state.get("ob_company_name", "").strip(),
            st.session_state.get("ob_sector", "").strip(),
            st.session_state.get("ob_industry_cluster", "").strip(),
        ],
        2: [
            len(st.session_state.get("ob_joint_committees", [])) > 0,
            st.session_state.get("ob_headquarters", "").strip(),
            st.session_state.get("ob_countries", "").strip(),
        ],
        3: [
            st.session_state.get("ob_operations", "").strip(),
            st.session_state.get("ob_employees_total", 0) > 0,
            st.session_state.get("ob_employees_belgium", 0) > 0,
        ],
        4: [
            len(st.session_state.get("ob_contract_types", [])) > 0,
            st.session_state.get("ob_payroll_frequency_blue", "").strip(),
            st.session_state.get("ob_payroll_frequency_white", "").strip(),
        ],
        5: [
            st.session_state.get("ob_weekly_hours", 0) > 0,
            st.session_state.get("ob_shift_work", "").strip(),
            st.session_state.get("ob_remote_policy", "").strip(),
        ],
        6: [
            st.session_state.get("ob_union_presence", "").strip(),
            st.session_state.get("ob_existing_policies", "").strip(),
            st.session_state.get("ob_priorities", "").strip(),
        ],
        7: [
            st.session_state.get("ob_industry_specific_details", "").strip(),
            st.session_state.get("ob_industry_specific_valid", False),
        ],
    }
    return all(checks.get(step, []))

def build_profile_from_onboarding_state():
    return {
        "_profile_completed": True,
        "company_name": st.session_state.get("ob_company_name", "").strip(),
        "sector": st.session_state.get("ob_sector", "").strip(),
        "joint_committees": ", ".join(st.session_state.get("ob_joint_committees", [])),
        "operations": st.session_state.get("ob_operations", "").strip(),
        "industry_cluster": st.session_state.get("ob_industry_cluster", "").strip(),
        "headquarters": st.session_state.get("ob_headquarters", "").strip(),
        "countries": st.session_state.get("ob_countries", "").strip(),
        "language": st.session_state.get("ob_language", "English"),
        "payroll_frequency_blue": st.session_state.get("ob_payroll_frequency_blue", "Weekly"),
        "payroll_frequency_white": st.session_state.get("ob_payroll_frequency_white", "Monthly"),
        "employees_total": int(st.session_state.get("ob_employees_total", 1)),
        "employees_belgium": int(st.session_state.get("ob_employees_belgium", 1)),
        "contract_types": ", ".join(st.session_state.get("ob_contract_types", [])),
        "weekly_hours": int(st.session_state.get("ob_weekly_hours", 38)),
        "shift_work": st.session_state.get("ob_shift_work", "No"),
        "remote_policy": st.session_state.get("ob_remote_policy", "No remote"),
        "union_presence": st.session_state.get("ob_union_presence", "Unknown"),
        "existing_policies": st.session_state.get("ob_existing_policies", "").strip(),
        "priorities": st.session_state.get("ob_priorities", "").strip(),
        "open_questions": st.session_state.get("ob_open_questions", "").strip(),
        "industry_specific_details": st.session_state.get("ob_industry_specific_details", "").strip(),
    }

def hydrate_onboarding_from_snapshot(company_id):
    """Pre-fill onboarding fields from saved company profile."""
    snapshot = get_company_profile_snapshot_cached(company_id)
    if not snapshot:
        return
    mapping = {
        "company_name": "ob_company_name",
        "sector": "ob_sector",
        "industry_cluster": "ob_industry_cluster",
        "operations": "ob_operations",
        "headquarters": "ob_headquarters",
        "countries": "ob_countries",
        "language": "ob_language",
        "payroll_frequency_blue": "ob_payroll_frequency_blue",
        "payroll_frequency_white": "ob_payroll_frequency_white",
        "employees_total": "ob_employees_total",
        "employees_belgium": "ob_employees_belgium",
        "weekly_hours": "ob_weekly_hours",
        "shift_work": "ob_shift_work",
        "remote_policy": "ob_remote_policy",
        "union_presence": "ob_union_presence",
        "existing_policies": "ob_existing_policies",
        "priorities": "ob_priorities",
        "open_questions": "ob_open_questions",
        "industry_specific_details": "ob_industry_specific_details",
    }
    for src, dst in mapping.items():
        if dst in st.session_state and str(st.session_state.get(dst)).strip():
            continue
        value = snapshot.get(src)
        if value is None or value == "":
            continue
        if dst in {"ob_employees_total", "ob_employees_belgium", "ob_weekly_hours"}:
            try:
                st.session_state[dst] = int(float(value))
            except Exception:
                pass
        else:
            st.session_state[dst] = value

    pcs = snapshot.get("joint_committees", "")
    if pcs and not st.session_state.get("ob_joint_committees"):
        st.session_state["ob_joint_committees"] = [p.strip() for p in pcs.split(",") if p.strip()]

    cts = snapshot.get("contract_types", "")
    if cts and not st.session_state.get("ob_contract_types"):
        st.session_state["ob_contract_types"] = [c.strip() for c in cts.split(",") if c.strip()]

def onboarding_page():
    hydrate_onboarding_from_snapshot(st.session_state.company_id)

    total_steps = 7
    step = st.session_state.get("onboarding_step", 1)
    step = max(1, min(total_steps, step))
    st.session_state.onboarding_step = step

    st.markdown("""
    <div class="onboarding-shell animated-panel">
        <h1 class="onboarding-title">Set Up Your Company Profile</h1>
        <p class="onboarding-subtitle">Guided setup in short steps. We ask a few questions at a time so Friday can apply Belgian labor law and CAO rules accurately.</p>
    </div>
    """, unsafe_allow_html=True)

    st.progress(step / total_steps, text=f"Step {step} of {total_steps}")
    st.caption("Maximum 3 questions per step.")

    if step == 1:
        st.markdown("### 1. Company Basics")
        st.text_input("Company name *", key="ob_company_name")
        st.text_input("Sector *", placeholder="Hospitality, Manufacturing, Retail...", key="ob_sector")
        st.selectbox(
            "Industry cluster *",
            ["Hospitality", "Manufacturing", "Retail", "Logistics/Transport", "Healthcare", "Professional Services", "Other"],
            key="ob_industry_cluster"
        )
    elif step == 2:
        st.markdown("### 2. Legal Scope")
        suggested_pcs = suggest_joint_committees(
            st.session_state.get("ob_sector", ""),
            st.session_state.get("ob_industry_cluster", "")
        )
        st.info(f"Suggested PCs based on sector: {', '.join(suggested_pcs)}")
        pc_options = sorted(set(COMMON_PC_OPTIONS + suggested_pcs))
        st.multiselect(
            "Applicable Joint Committees (PC) *",
            options=pc_options,
            default=[pc for pc in st.session_state.get("ob_joint_committees", []) if pc in pc_options],
            key="ob_joint_committees"
        )
        st.text_input("Headquarters (city, country) *", placeholder="Brussels, Belgium", key="ob_headquarters")
        st.text_input("Countries of operation *", value=st.session_state.get("ob_countries", "Belgium"), key="ob_countries")
    elif step == 3:
        st.markdown("### 3. Workforce Snapshot")
        st.text_area("Main business activities *", height=100, key="ob_operations")
        st.number_input("Total employees *", min_value=1, step=1, key="ob_employees_total")
        st.number_input("Employees in Belgium *", min_value=1, step=1, key="ob_employees_belgium")
    elif step == 4:
        st.markdown("### 4. Contracts and Payroll")
        st.multiselect(
            "Contract types used *",
            ["Full-time", "Part-time", "Fixed-term", "Temporary agency", "Student", "Freelancers"],
            key="ob_contract_types"
        )
        st.selectbox("Blue-collar payroll frequency *", ["Weekly", "Every 2 weeks", "Monthly"], key="ob_payroll_frequency_blue")
        st.selectbox("White-collar payroll frequency *", ["Monthly", "Every 2 weeks", "Weekly"], key="ob_payroll_frequency_white")
    elif step == 5:
        st.markdown("### 5. Work Pattern")
        st.number_input("Standard weekly hours *", min_value=1, max_value=60, value=38, step=1, key="ob_weekly_hours")
        st.selectbox("Shift/night/weekend work *", ["No", "Occasionally", "Regularly"], key="ob_shift_work")
        st.selectbox("Remote work policy *", ["No remote", "Hybrid", "Fully remote"], key="ob_remote_policy")
    elif step == 6:
        st.markdown("### 6. Policies and Priorities")
        st.selectbox("Union delegation present *", ["Yes", "No", "Unknown"], key="ob_union_presence")
        st.text_area("Current HR policy summary *", height=110, key="ob_existing_policies")
        st.text_area(
            "Top compliance priorities *",
            height=90,
            placeholder="Working time, overtime, leave, dismissal, wage indexation...",
            key="ob_priorities"
        )
        st.text_area("Open HR/legal questions (optional)", height=80, key="ob_open_questions")
    else:
        st.markdown("### 7. Industry Details")
        details, is_valid = get_industry_specific_payload(st.session_state.get("ob_industry_cluster", "Other"))
        st.session_state.ob_industry_specific_details = details
        st.session_state.ob_industry_specific_valid = is_valid
        st.file_uploader(
            "Upload company policy documents (optional)",
            type=["pdf", "docx", "xlsx", "pptx"],
            accept_multiple_files=True,
            key="ob_policy_files"
        )

    col_back, col_next = st.columns([1, 1])
    with col_back:
        if step > 1 and st.button("Back", use_container_width=True):
            st.session_state.onboarding_step = step - 1
            st.rerun()
    with col_next:
        if step < total_steps:
            if st.button("Next", use_container_width=True, type="primary"):
                if not validate_onboarding_step(step):
                    st.error("Please complete the required fields on this step.")
                else:
                    st.session_state.onboarding_step = step + 1
                    st.rerun()
        else:
            if st.button("Save Company Profile", use_container_width=True, type="primary"):
                if not validate_onboarding_step(step):
                    st.error("Please complete the required fields on this step.")
                else:
                    profile = build_profile_from_onboarding_state()
                    ok = index_company_profile(st.session_state.company_id, profile)
                    if ok:
                        policy_files = st.session_state.get("ob_policy_files") or []
                        upload_results = {"success": 0, "exists": 0, "error": 0}
                        for f in policy_files:
                            res = process_and_store_document(f, st.session_state.company_id, force_overwrite=False)
                            if res == "success":
                                upload_results["success"] += 1
                            elif res == "exists":
                                upload_results["exists"] += 1
                            else:
                                upload_results["error"] += 1

                        if upload_results["success"] > 0:
                            st.info(f"Policy uploads indexed: {upload_results['success']}.")
                        if upload_results["error"] > 0:
                            st.warning(f"Some policy uploads failed: {upload_results['error']}.")

                        st.success("Company profile saved. You can now use FRIDAY.")
                        st.session_state.onboarding_required = False
                        st.session_state.onboarding_step = 1
                        st.session_state.show_welcome_anim = True
                        st.session_state.view = "chat"
                        if not st.session_state.current_chat_id:
                            create_new_chat()
                        time.sleep(0.8)
                        st.rerun()
                    else:
                        st.error("Could not save the company profile. Please try again.")


# --- 5. AUTHENTICATION ---
def handle_login():
    """Callback for login form - handles authentication state properly."""
    pw = st.session_state.get("login_password", "")
    if not pw:
        st.session_state.login_error = "Please enter an access code"
        return
    try:
        res = supabase.table('clients').select("*").eq('access_code', pw).execute()
        if res.data:
            st.session_state.authenticated = True
            st.session_state.company_id = res.data[0]['company_id']
            st.session_state.onboarding_required = not company_profile_exists(st.session_state.company_id)
            st.session_state.view = "onboarding" if st.session_state.onboarding_required else "chat"
            st.session_state.onboarding_step = 1
            st.session_state.login_error = None
        else: st.session_state.login_error = "Invalid Code"
    except: st.session_state.login_error = "Connection Error"

def login_page():
    if "login_error" not in st.session_state: st.session_state.login_error = None
    st.markdown("<br>", unsafe_allow_html=True)
    col1, col2, col3 = st.columns([1, 2, 1])
    with col2:
        st.markdown("""
        <div style="text-align: center; margin-bottom: 48px;">
            <div style="font-family: 'Playfair Display', serif; font-size: 72px; font-weight: 800; color: #1A3C34;">Friday</div>
            <h1 class="hero-title" style="font-size: 42px;">Your Intelligent<br><em>HR Companion</em></h1>
        </div>
        """, unsafe_allow_html=True)
        
        pw = st.text_input("Access Code", type="password", key="login_password")
        st.markdown("<div style='height: 16px;'></div>", unsafe_allow_html=True)
        if st.button("Sign In", use_container_width=True, type="primary", on_click=handle_login):
            if st.session_state.authenticated: st.rerun()
        
        if st.session_state.login_error: st.error(st.session_state.login_error)

if not st.session_state.authenticated: login_page()
else:
    if st.session_state.onboarding_required:
        onboarding_page()
    else:
        render_sidebar()
        if st.session_state.view == "chat": chat_page()
        elif st.session_state.view == "documents": documents_page()
